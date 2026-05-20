from pathlib import Path
import sys


sys.path.insert(0, str(Path(__file__).parent / ".deps_pptx"))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path("outputs") / "ppt_shapes"
OUT_FILE = OUT_DIR / "mcp_interface_architecture_shapes.pptx"


FONT = "Malgun Gothic"


def c(hex_color: str) -> RGBColor:
    hex_color = hex_color.strip("#")
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16),
    )


def inch(v: float):
    return Inches(v)


def set_fill(shape, color=None):
    if color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = c(color)


def set_line(shape, color="#BFBFBF", width=0.75):
    if color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = c(color)
        shape.line.width = Pt(width)


def rect(slide, x, y, w, h, fill="#FFFFFF", line="#BFBFBF", width=0.75, radius=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        inch(x),
        inch(y),
        inch(w),
        inch(h),
    )
    set_fill(shp, fill)
    set_line(shp, line, width)
    return shp


def oval(slide, x, y, w, h, fill="#FFFFFF", line="#BFBFBF", width=0.75):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, inch(x), inch(y), inch(w), inch(h))
    set_fill(shp, fill)
    set_line(shp, line, width)
    return shp


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    size=12,
    color="#333333",
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    fill=None,
    line=None,
    margin=0.03,
    rotation=0,
):
    box = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    box.rotation = rotation
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = inch(margin)
    tf.margin_right = inch(margin)
    tf.margin_top = inch(0.015)
    tf.margin_bottom = inch(0.015)
    tf.vertical_anchor = valign
    if fill is not None:
        set_fill(box, fill)
    if line is not None:
        set_line(box, line)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = c(color)
    return box


def add_multiline_text(slide, x, y, w, h, lines, size=9, color="#333333", bold_first=False):
    box = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = inch(0.03)
    tf.margin_right = inch(0.03)
    tf.margin_top = inch(0.02)
    tf.margin_bottom = inch(0.02)
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(0)
        run = p.add_run()
        run.text = line
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bool(bold_first and idx == 0)
        run.font.color.rgb = c(color)
    return box


def connector(slide, x1, y1, x2, y2, color="#777777", width=1.0):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, inch(x1), inch(y1), inch(x2), inch(y2)
    )
    line.line.color.rgb = c(color)
    line.line.width = Pt(width)
    return line


def arrow_down(slide, x, y1, y2, color="#333333", width=1.0, tri=0.08):
    connector(slide, x, y1, x, y2 - tri * 0.5, color, width)
    head = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE, inch(x - tri / 2), inch(y2 - tri), inch(tri), inch(tri)
    )
    head.rotation = 180
    set_fill(head, color)
    set_line(head, None)
    return head


def arrow_up(slide, x, y1, y2, color="#333333", width=1.0, tri=0.08):
    connector(slide, x, y1 + tri * 0.5, x, y2, color, width)
    head = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE, inch(x - tri / 2), inch(y1), inch(tri), inch(tri)
    )
    set_fill(head, color)
    set_line(head, None)
    return head


def robot(slide, cx, cy, scale=1.0, color="#2D2D2D"):
    w = 0.23 * scale
    h = 0.18 * scale
    head = rect(slide, cx - w / 2, cy - h / 2, w, h, fill=color, line=color, radius=True)
    oval(slide, cx - w * 0.28, cy - h * 0.08, w * 0.16, h * 0.16, fill="#FFFFFF", line="#FFFFFF")
    oval(slide, cx + w * 0.12, cy - h * 0.08, w * 0.16, h * 0.16, fill="#FFFFFF", line="#FFFFFF")
    connector(slide, cx, cy - h / 2, cx, cy - h / 2 - 0.06 * scale, color=color, width=1.0)
    oval(slide, cx - 0.025 * scale, cy - h / 2 - 0.095 * scale, 0.05 * scale, 0.05 * scale, fill=color, line=color)
    rect(slide, cx - 0.11 * scale, cy + h / 2 + 0.02 * scale, 0.22 * scale, 0.10 * scale, fill=color, line=color, radius=True)
    connector(slide, cx - 0.15 * scale, cy + 0.02 * scale, cx - 0.22 * scale, cy + 0.12 * scale, color=color, width=1.0)
    connector(slide, cx + 0.15 * scale, cy + 0.02 * scale, cx + 0.22 * scale, cy + 0.12 * scale, color=color, width=1.0)


def database_icon(slide, x, y, w, h, color="#2E2E2E"):
    top = oval(slide, x, y, w, h * 0.32, fill="#FFFFFF", line=color, width=1.2)
    body = rect(slide, x, y + h * 0.14, w, h * 0.62, fill="#FFFFFF", line=color, width=1.2)
    bottom = oval(slide, x, y + h * 0.58, w, h * 0.32, fill="#FFFFFF", line=color, width=1.2)
    top.fill.transparency = 100000
    body.fill.transparency = 100000
    bottom.fill.transparency = 100000


def small_diagonal_stripes(slide, x, y, w, h):
    for i in range(-3, int(w / 0.08) + 8):
        sx = x + i * 0.08
        connector(slide, sx, y + h, sx + 0.22, y, color="#D9D9D9", width=0.35)


def add_feature_box(slide, x, y, w, h, title, subtitle, icon_kind):
    rect(slide, x, y, w, h, fill="#FFFFFF", line="#C6C6C6", width=1.1, radius=True)
    ix = x + 0.13
    iy = y + 0.18
    if icon_kind == "cube":
        cube1 = slide.shapes.add_shape(MSO_SHAPE.CUBE, inch(ix), inch(iy), inch(0.22), inch(0.22))
        set_fill(cube1, "#FFFFFF")
        set_line(cube1, "#333333", 1.1)
        cube2 = slide.shapes.add_shape(
            MSO_SHAPE.CUBE, inch(ix + 0.18), inch(iy + 0.13), inch(0.18), inch(0.18)
        )
        set_fill(cube2, "#FFFFFF")
        set_line(cube2, "#333333", 1.1)
    elif icon_kind == "shield":
        shp = slide.shapes.add_shape(MSO_SHAPE.PENTAGON, inch(ix), inch(iy - 0.02), inch(0.30), inch(0.34))
        shp.rotation = 180
        set_fill(shp, "#FFFFFF")
        set_line(shp, "#333333", 1.3)
    elif icon_kind == "gear":
        shp = slide.shapes.add_shape(MSO_SHAPE.GEAR_6, inch(ix), inch(iy - 0.02), inch(0.33), inch(0.33))
        set_fill(shp, "#FFFFFF")
        set_line(shp, "#333333", 1.2)
    elif icon_kind == "db":
        database_icon(slide, ix, iy - 0.02, 0.32, 0.34, color="#333333")
    add_text(slide, x + 0.56, y + 0.13, w - 0.63, 0.18, title, size=6.6, color="#222222", bold=True)
    add_text(slide, x + 0.56, y + 0.32, w - 0.63, 0.22, subtitle, size=5.6, color="#333333")


def add_adapter(slide, x, y, w, h, title, subtitle, icon_kind):
    rect(slide, x, y, w, h, fill="#FFFFFF", line="#D0D0D0", width=1.1, radius=True)
    ix = x + 0.16
    iy = y + 0.15
    if icon_kind == "globe":
        oval(slide, ix, iy, 0.29, 0.29, fill="#FFFFFF", line="#333333", width=1.2)
        connector(slide, ix + 0.145, iy, ix + 0.145, iy + 0.29, color="#333333", width=0.9)
        connector(slide, ix, iy + 0.145, ix + 0.29, iy + 0.145, color="#333333", width=0.9)
        connector(slide, ix + 0.04, iy + 0.06, ix + 0.25, iy + 0.06, color="#333333", width=0.7)
        connector(slide, ix + 0.04, iy + 0.23, ix + 0.25, iy + 0.23, color="#333333", width=0.7)
    elif icon_kind == "nodes":
        for dx, dy in [(0.04, 0.02), (0.27, 0.15), (0.04, 0.29)]:
            oval(slide, ix + dx, iy + dy, 0.055, 0.055, fill="#FFFFFF", line="#333333", width=1.0)
        connector(slide, ix + 0.07, iy + 0.06, ix + 0.28, iy + 0.18, color="#333333", width=1.0)
        connector(slide, ix + 0.07, iy + 0.31, ix + 0.28, iy + 0.18, color="#333333", width=1.0)
    elif icon_kind == "cloud":
        shp = slide.shapes.add_shape(MSO_SHAPE.CLOUD, inch(ix), inch(iy + 0.03), inch(0.35), inch(0.24))
        set_fill(shp, "#FFFFFF")
        set_line(shp, "#333333", 1.2)
    add_text(slide, x + 0.58, y + 0.12, w - 0.65, 0.18, title, size=7.2, color="#222222", bold=True)
    add_text(slide, x + 0.58, y + 0.31, w - 0.65, 0.18, subtitle, size=6.0, color="#333333")


def add_right_item(slide, idx, y, title, bullets):
    x = 9.53
    w = 2.98
    h = 0.34
    rect(slide, x, y, w, h, fill="#ECECEC", line=None)
    small_diagonal_stripes(slide, x, y, w, h)
    rect(slide, x + 0.04, y + 0.045, 0.25, 0.25, fill="#666159", line=None)
    add_text(
        slide,
        x + 0.04,
        y + 0.057,
        0.25,
        0.20,
        str(idx),
        size=10,
        color="#FFFFFF",
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        bold=False,
    )
    add_text(slide, x + 0.35, y + 0.045, 2.50, 0.27, title, size=12.3, color="#444444", bold=True)
    add_multiline_text(slide, x + 0.04, y + 0.42, 2.82, 0.58, [f"• {b}" for b in bullets], size=9.1, color="#222222")


def main():
    prs = Presentation()
    prs.slide_width = inch(13)
    prs.slide_height = inch(9)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Page background and header.
    rect(slide, 0, 0, 13, 9, fill="#FFFFFF", line=None)
    rect(slide, 0, 0, 13, 0.09, fill="#746D61", line=None)
    rect(slide, 0, 0.09, 13, 1.78, fill="#EFEFEF", line=None)

    add_text(slide, 0.08, 0.16, 0.55, 0.55, "03", size=31, color="#9B9B9B")
    add_text(slide, 0.71, 0.37, 2.6, 0.25, "기술제안-공통부문", size=12.5, color="#333333")
    add_text(slide, 0.39, 0.82, 3.9, 0.28, "3.2 멀티 에이전트 시스템 설계", size=12.4, color="#444444")
    add_text(slide, 0.50, 1.20, 4.5, 0.33, "3) Tool Interface (MCP Server)", size=15.0, color="#111111")

    for i, ch in enumerate("수행방안"):
        oval(slide, 11.36 + i * 0.36, 0.16, 0.33, 0.33, fill="#C89400", line="#C89400")
        add_text(slide, 11.36 + i * 0.36, 0.20, 0.33, 0.20, ch, size=12.0, color="#FFFFFF", bold=True, align=PP_ALIGN.CENTER)

    # Left-side cropped navigation circles visible in the source image.
    oval(slide, -1.35, 2.47, 1.86, 1.86, fill="#F37320", line="#AE4E17", width=1.0)
    add_text(slide, -0.11, 2.95, 0.30, 1.00, "구현\n방안", size=14, color="#FFFFFF", bold=True, align=PP_ALIGN.CENTER, rotation=270)
    oval(slide, -1.31, 4.47, 1.86, 1.86, fill="#E9E9E9", line="#979797", width=0.9)

    # Section headers.
    rect(slide, 0.35, 2.20, 8.90, 0.35, fill="#6E675A", line=None)
    add_text(slide, 0.35, 2.275, 8.90, 0.18, "MCP Interface 아키텍처", size=13.2, color="#FFFFFF", align=PP_ALIGN.CENTER)
    rect(slide, 9.42, 2.20, 3.18, 0.35, fill="#6E675A", line=None)
    add_text(slide, 9.42, 2.27, 3.18, 0.20, "구축 전략 / 핵심 기능", size=13.2, color="#FFFFFF", align=PP_ALIGN.CENTER)

    rect(slide, 0.35, 2.55, 8.90, 5.59, fill=None, line="#D0D0D0", width=0.8)
    rect(slide, 9.49, 2.64, 3.10, 5.50, fill=None, line="#D0D0D0", width=0.8)

    # Main architecture diagram.
    rect(slide, 2.11, 2.87, 4.12, 0.61, fill="#FFFFFF", line="#B8B8B8", width=1.1, radius=True)
    robot(slide, 2.57, 3.19, 1.28)
    robot(slide, 3.03, 3.19, 1.28)
    robot(slide, 3.49, 3.19, 1.28)
    add_text(slide, 3.97, 3.02, 1.85, 0.22, "AI Agent / Multi-Agent", size=10.8, color="#111111", bold=True)
    add_text(slide, 3.97, 3.25, 1.95, 0.18, "(Planner / Analyzer / Executor 등)", size=7.4, color="#111111")

    arrow_down(slide, 4.12, 3.49, 3.67, color="#222222", width=1.0)
    rect(slide, 2.11, 3.68, 4.14, 0.45, fill="#FFFFFF", line="#B8B8B8", width=1.1, radius=True)
    add_text(slide, 2.11, 3.73, 4.14, 0.17, "MCP Client", size=9.5, color="#333333", bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, 2.11, 3.93, 4.14, 0.17, "Tool Discovery / 호출 / 결과 처리", size=7.6, color="#333333", align=PP_ALIGN.CENTER)
    arrow_down(slide, 4.12, 4.13, 4.29, color="#222222", width=1.0)

    rect(slide, 0.63, 4.30, 6.33, 1.57, fill="#FFFFFF", line="#444444", width=1.3, radius=True)
    rect(slide, 0.64, 4.32, 6.30, 0.33, fill="#5C4124", line=None)
    add_text(slide, 0.64, 4.39, 6.30, 0.18, "MCP Server (Tool Hub)", size=11.0, color="#FFFFFF", align=PP_ALIGN.CENTER)
    add_text(slide, 0.64, 4.70, 6.30, 0.18, "레거시 시스템 표준 추상화 계층", size=8.5, color="#222222", bold=True, align=PP_ALIGN.CENTER)

    add_feature_box(slide, 0.70, 5.00, 1.55, 0.62, "Tool 등록/관리", "Tool 메타데이터 관리 및\n버전 관리", "cube")
    add_feature_box(slide, 2.33, 5.00, 1.50, 0.62, "인증/권한 관리", "보안 인증, 접근 제어 및\n권한 관리", "shield")
    add_feature_box(slide, 3.90, 5.00, 1.50, 0.62, "프로토콜 변환", "요청/응답 표준화 및\n데이터 포맷 변환", "gear")
    add_feature_box(slide, 5.46, 5.00, 1.43, 0.62, "모니터링/로깅", "호출 이력, 성능 모니터링 및\n감사 로그 관리", "db")

    for x in [1.66, 3.84, 5.99]:
        arrow_up(slide, x, 5.87, 6.15, color="#444444", width=1.0, tri=0.075)

    add_adapter(slide, 0.62, 6.15, 2.00, 0.55, "REST API Adapter", "HTTP/HTTPS 기반 연동", "globe")
    add_adapter(slide, 2.91, 6.15, 1.87, 0.55, "gRPC Adapter", "gRPC 기반 고성능 연동", "nodes")
    add_adapter(slide, 5.02, 6.15, 2.00, 0.55, "External API Adapter", "SaaS / 외부 API 연동", "cloud")

    for x in [1.66, 3.84, 5.99]:
        connector(slide, x, 6.70, x, 7.02, color="#9A9A9A", width=1.2)
        arrow_down(slide, x, 6.71, 6.82, color="#9A9A9A", width=1.0, tri=0.06)
        arrow_up(slide, x, 6.91, 7.02, color="#9A9A9A", width=1.0, tri=0.06)

    rect(slide, 0.62, 7.02, 6.34, 0.74, fill="#FFFFFF", line="#D0D0D0", width=1.1, radius=True)
    rect(slide, 0.88, 7.26, 0.33, 0.13, fill="#333333", line="#333333", radius=True)
    rect(slide, 0.88, 7.43, 0.33, 0.13, fill="#333333", line="#333333", radius=True)
    rect(slide, 0.88, 7.60, 0.33, 0.13, fill="#333333", line="#333333", radius=True)
    for yy in [7.29, 7.46, 7.63]:
        oval(slide, 0.93, yy, 0.03, 0.03, fill="#FFFFFF", line="#FFFFFF")
    add_text(slide, 1.60, 7.20, 2.45, 0.20, "Legacy / Internal Systems", size=11.5, color="#222222", bold=True)
    add_text(slide, 1.70, 7.50, 0.70, 0.16, "Core System", size=6.6, color="#333333")
    for x in [2.58, 3.41, 4.24, 5.06, 6.29]:
        connector(slide, x, 7.45, x, 7.62, color="#C8C8C8", width=0.8)
    add_text(slide, 2.75, 7.50, 0.33, 0.16, "ERP", size=6.6, color="#333333")
    add_text(slide, 3.72, 7.50, 0.33, 0.16, "CRM", size=6.6, color="#333333")
    add_text(slide, 4.47, 7.50, 0.55, 0.16, "DW / DB", size=6.6, color="#333333")
    add_text(slide, 5.45, 7.50, 0.80, 0.16, "File / Batch", size=6.6, color="#333333")
    add_text(slide, 6.60, 7.50, 0.25, 0.16, "...", size=7.0, color="#333333")

    rect(slide, 7.20, 4.32, 1.87, 3.38, fill="#E9E9E9", line="#A8A8A8", width=0.8, radius=True)
    add_text(slide, 7.35, 4.52, 1.55, 0.22, "MCP 표준구조", size=12.0, color="#444444", bold=True, align=PP_ALIGN.CENTER)
    add_multiline_text(
        slide,
        7.34,
        4.95,
        1.55,
        2.45,
        [
            "Tool",
            "• 기능을 표준 Tool 정의\n  (입력/출력/설명/버전)",
            "Tool Schema",
            "• 일관된 스키마 기반\n  입출력 정의",
            "Transport",
            "• 표준 전송 프로토콜",
            "  (stdio/http+SSE/Streamable HTTP)",
            "Response",
            "• 표준 응답 포맷",
        ],
        size=7.7,
        color="#3E3E3E",
        bold_first=True,
    )

    # Right strategy / key features panel.
    add_right_item(slide, 1, 2.64, "레거시 표준 추상화 계층", ["레거시 API를 Tool로 추상화", "MCP 서버가 표준 인터페이스 제공", "시스템 종속성 최소화"])
    add_right_item(slide, 2, 3.91, "MCP 기반 Tool 관리", ["Tool 등록/관리 및 버전 관리", "인증/권한 및 보안 관리"])
    add_right_item(slide, 3, 4.93, "다양한 프로토콜 지원", ["REST / gRPC / External API 지원", "프로토콜 변환 및 데이터 포맷 변환"])
    add_right_item(slide, 4, 6.01, "Agent 일관 호출 인터페이스", ["MCP Client를 통한 Tool 호출", "동일한 방식으로 모든 Tool 접근"])
    add_right_item(slide, 5, 7.10, "Tool 추가 / 변경", ["새로운 Tool 추가 변경 용이", "비침투 방식으로 유지보수 비용 절감"])
    add_text(slide, 9.50, 8.23, 3.05, 0.13, "* 비침투 방식:기존 핵심 코드 수정없이 유지관리", size=6.5, color="#333333")

    # Bottom logos are recreated as editable text/shape marks.
    add_text(slide, 0.19, 8.72, 0.25, 0.22, "✳", size=16, color="#F9A602")
    add_text(slide, 0.39, 8.73, 1.05, 0.18, "KB국민카드", size=9.2, color="#777777", bold=True)
    rect(slide, 10.98, 8.72, 0.21, 0.25, fill="#23613F", line=None)
    add_text(slide, 11.005, 8.765, 0.16, 0.12, "솔", size=6.3, color="#FFFFFF", bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, 11.28, 8.70, 0.78, 0.16, "Pinetree", size=11.0, color="#777777", bold=True)
    add_text(slide, 11.29, 8.86, 0.62, 0.08, "partners", size=3.6, color="#777777")
    add_text(slide, 12.08, 8.71, 0.65, 0.16, "QuinTet", size=11.0, color="#B13A35", bold=True)
    add_text(slide, 12.25, 8.87, 0.45, 0.07, "The World Class ICT Provider", size=2.7, color="#777777")

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_FILE)
    print(OUT_FILE.resolve())


if __name__ == "__main__":
    main()
